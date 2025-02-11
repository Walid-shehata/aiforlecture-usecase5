import streamlit as st
import boto3
import os
import json
from dotenv import load_dotenv
from subjects import get_subjects
from chapters import get_chapters
from topicSummaryCreator import get_topics
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import re
import logging
import traceback
logging.basicConfig(level=logging.DEBUG)

load_dotenv()

# Initialize AWS clients
s3 = boto3.client('s3',
                  aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
                  aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY'),
                  region_name=os.getenv('AWS_REGION')
                  )

bedrock_runtime = boto3.client('bedrock-runtime',
                               aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
                               aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY'),
                               region_name=os.getenv('AWS_REGION')
                               )
bedrock_agent_runtime = boto3.client(service_name='bedrock-agent-runtime', region_name=os.getenv('AWS_REGION'))

BUCKET_NAME = os.getenv('S3_BUCKET_NAME')
ARIFACTS_BUCKET_NAME = os.getenv('S3_ARTIFACTS_BUCKET_NAME')
KNOWLEDGE_BASE_ID = os.getenv('BEDROCK_KNOWLEDGE_BASE_ID')

def generate_bulleted_content(content):
    prompt = f"Based on the following content, generate 3-4 concise bullet points that summarize the key ideas:\n\n{content}"

    response = bedrock_runtime.invoke_model(
        modelId="anthropic.claude-3-sonnet-20240229-v1:0",
        contentType="application/json",
        accept="application/json",
        body=json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 300,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.3,
            "top_p": 1.0,
        })
    )
    response_body = json.loads(response['body'].read())
    return response_body['content'][0]['text'].strip()

def generate_conclusion_summary(structure):
    content = "\n".join([slide['content'] for slide in structure if slide['type'] in ['Title&Text', 'Other']])
    prompt = f"Based on the following content from the presentation, generate a concise conclusion summary with 3-4 bullet points:\n\n{content}"

    response = bedrock_runtime.invoke_model(
        modelId="anthropic.claude-3-sonnet-20240229-v1:0",
        contentType="application/json",
        accept="application/json",
        body=json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 300,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.3,
            "top_p": 1.0,
        })
    )
    response_body = json.loads(response['body'].read())
    return response_body['content'][0]['text'].strip()




def parse_presentation_structure(structure, subject, chapter, selected_topics):
    slides = []
    for slide in structure.split('\n'):
        slide = slide.strip()
        if not slide:
            continue

        match = re.match(r'(?:Slide\s*)?(\d+)[\.,]\s*(\w+(?:&\w+)?),\s*(.+)', slide)
        if match:
            number, slide_type, content = match.groups()
            if slide_type in ['Poll', 'Discussion']:
                slide_type = 'Other'

            if slide_type == 'TitleOnly':
                title = content
                content = ''
            else:
                title, *content_parts = content.split(',', 1)
                content = content_parts[0].strip() if content_parts else ''

            if number == '1' and slide_type == 'TitleOnly':
                title = f"Introduction: {subject} - {chapter}"
                content = f"Topics: {', '.join(selected_topics)}"

            slides.append({
                'number': int(number),
                'type': slide_type,
                'title': title.strip(),
                'content': content,
                'image_prompt': content if slide_type == 'Title&Picture' else ''
            })

    st.write(f"Parsed {len(slides)} slides")  # Debug output
    return slides


def generate_slide_notes(slide_content):
    prompt = f"Generate detailed speaker notes for the following slide content:\n\n{slide_content}"

    response = bedrock_runtime.invoke_model(
        modelId="anthropic.claude-3-sonnet-20240229-v1:0",
        contentType="application/json",
        accept="application/json",
        body=json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 500,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.3,
            "top_p": 1.0,
        })
    )
    response_body = json.loads(response['body'].read())
    return response_body['content'][0]['text'].strip()

def generate_image(prompt):
    try:
        response = bedrock_runtime.invoke_model(
            modelId="stability.stable-diffusion-xl-v1",
            contentType="application/json",
            accept="image/png",
            body=json.dumps({
                "text_prompts": [{"text": prompt}],
                "cfg_scale": 10,
                "steps": 50,
                "seed": 42,
            })
        )
        return response['body'].read()
    except Exception as e:
        st.error(f"Error generating image: {str(e)}")
        return None


def create_powerpoint(structure):
    template_path = "Anyuniversity.pptx"
    try:
        prs = Presentation(template_path)
    except Exception as e:
        st.error(f"Failed to load template: {str(e)}")
        return None

    progress_bar = st.progress(0)
    total_slides = len(structure)

    for i, slide in enumerate(structure):
        slide_type = slide['type']
        if slide_type == 'TitleOnly':
            slide_layout = prs.slide_layouts[0]
        elif slide_type == 'Title&Text':
            slide_layout = prs.slide_layouts[1]
        elif slide_type == 'Title&Picture':
            slide_layout = prs.slide_layouts[8]
        else:
            slide_layout = prs.slide_layouts[2]

        slide_obj = prs.slides.add_slide(slide_layout)

        # Add title
        if slide_obj.shapes.title:
            slide_obj.shapes.title.text = slide['title']
        else:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            txBox = slide_obj.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = slide['title']

        if slide_type in ['Title&Text', 'Other']:
            bulleted_content = generate_bulleted_content(slide['content'])
            try:
                body_shape = slide_obj.placeholders[1]
                body_shape.text = bulleted_content
            except IndexError:
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5)
                txBox = slide_obj.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = slide['content']
        elif slide_type == 'Title&Picture':
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(5.5)
            txBox = slide_obj.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = f"[Suggested image: {slide['image_prompt']}]"

        notes_slide = slide_obj.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = generate_slide_notes(slide['content'])

        progress_bar.progress((i + 1) / total_slides)

    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    progress_bar.empty()
    return pptx_buffer


def review_and_edit_structure(structure):
    edited_structure = []
    slide_types = ['TitleOnly', 'Title&Text', 'Title&Picture', 'Other']

    for i, slide in enumerate(structure):
        with st.expander(f"Slide {slide['number']}: {slide['title']}"):
            delete_slide = st.checkbox("Delete Slide", key=f"delete_{i}")

            if not delete_slide:
                edited_type = st.selectbox(
                    "Slide Type",
                    slide_types,
                    key=f"type_{i}",
                    index=slide_types.index(slide['type'])
                )

                edited_title = st.text_input("Slide Title", slide['title'], key=f"title_{i}")

                if edited_type in ['Title&Text', 'Other']:
                    edited_content = st.text_area("Slide Content", slide['content'], key=f"content_{i}")
                elif edited_type == 'Title&Picture':
                    edited_content = st.text_area("Image Prompt", slide['image_prompt'], key=f"image_prompt_{i}")
                else:
                    edited_content = ''

                edited_structure.append({
                    'number': len(edited_structure) + 1,
                    'type': edited_type,
                    'title': edited_title,
                    'content': edited_content,
                    'image_prompt': edited_content if edited_type == 'Title&Picture' else ''
                })

    return edited_structure



def generate_presentation_structure(subject, chapter, selected_topics, lecture_length):
    num_slides = lecture_length // 3

    query = f"""Retrieve information for a presentation on:
    Subject: {subject}
    Chapter: {chapter}
    Topics: {', '.join(selected_topics)}
    """
    agent_response = bedrock_agent_runtime.retrieve(
        knowledgeBaseId=KNOWLEDGE_BASE_ID,
        retrievalQuery={'text': query},
        retrievalConfiguration={'vectorSearchConfiguration': {'numberOfResults': 6}}
    )

    context = "Based on the following information:\n\n"
    for result in agent_response['retrievalResults']:
        if 'text' in result['content']:
            context += f"- {result['content']['text']}\n"
        elif 'byteContent' in result['content']:
            content_type = result['content']['byteContent'].split(';')[0].split(':')[1]
            context += f"- [Content of type: {content_type}]\n"

    prompt = f"""{context}
    As an AI assistant for teachers, create a detailed structure for a PowerPoint presentation on the following:

    Subject: {subject}
    Chapter: {chapter}
    Topics: {', '.join(selected_topics)}
    Number of slides: {num_slides}
    The slides can be of the following Types: Title&Text, Title&Picture, TitleOnly(For Intro, separator slides, assignments, conclusion and Thank you)
    For Title&Text, return Title Name, and what will you discuss in the slides.
    For Title&Picture, return Title Name, and a text prompt that can be fed for an AI tool to generate a matching picture.
    For TitleOnly, return Title Name
    Provide a presentation structure in this format
    Slide Number, Slide Type, Returned Content(This will vary based on the Slide Type).
    The flow of the presentation should mimic story telling style to keep students engaged.
    To increase engagements, we need to have throughout the presentation discussion questions, Polls.
    Limit duscussion or polls to only two or three slides maximum.
    Start with an Intro, end with a Conclusion to summarize then a "Thank You" slide.
    Ensure that the total number of slides matches the specified {num_slides}.
    Use mainly the provided context, you can use your own Knowledge but only in limited situations.
    Give me the output I asked for in my format without any extra comment from you about it.
    """

    response = bedrock_runtime.invoke_model(
        modelId="anthropic.claude-3-sonnet-20240229-v1:0",
        contentType="application/json",
        accept="application/json",
        body=json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 2500,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.3,
            "top_p": 1.0,
        })
    )
    response_body = json.loads(response['body'].read())
    return response_body['content'][0]['text'].strip()

def lecture_planner():
    with st.expander("📚 Click here for Tool Instructions"):
        st.markdown("""
        **How to use the Lecture-Planner:**
        1. Select a subject from the first dropdown menu.
        2. Select a chapter from the second dropdown menu.
        3. Choose the topics you want to cover in your lecture using the multi-select box.
        4. Set the lecture length in minutes.
        5. Click "Generate Presentation Structure" to create an initial outline.
        6. Review the generated structure in the text area.
        7. Use the "Review and Edit Presentation Structure" section to modify individual slides:
           - You can change slide types, titles, and content.
           - Use the delete checkbox to remove unwanted slides.
        8. Click "Update Structure" to save your changes.
        9. When satisfied with the structure, click "Create PowerPoint Presentation".
        10. Download the generated PowerPoint file using the "Download PowerPoint Presentation" button.

        Note: 
        - The AI generates content based on the selected subject, chapter, and topics.
        - You can iterate through steps 7-9 to refine your presentation structure.
        - The final PowerPoint includes speaker notes for each slide.
        - Ensure all changes are saved before creating the PowerPoint.
         """)

    if 'subject' not in st.session_state:
        st.session_state.subject = None
    if 'chapter' not in st.session_state:
        st.session_state.chapter = None
    if 'selected_topics' not in st.session_state:
        st.session_state.selected_topics = []
    if 'lecture_length' not in st.session_state:
        st.session_state.lecture_length = 60
    if 'structure' not in st.session_state:
        st.session_state.structure = None
    if 'parsed_structure' not in st.session_state:
        st.session_state.parsed_structure = None

    subjects = get_subjects()
    st.session_state.subject = st.selectbox("Select Subject", subjects, key="LECTUREPLANNERSubjectSelector",
                                            index=subjects.index(st.session_state.subject) if st.session_state.subject in subjects else 0)

    if st.session_state.subject:
        chapters = get_chapters(st.session_state.subject)
        st.session_state.chapter = st.selectbox("Select Chapter", chapters, key="LECTUREPLANNERChapterSelector",
                                                index=chapters.index(st.session_state.chapter) if st.session_state.chapter in chapters else 0)

        if st.session_state.chapter:
            topics = get_topics(st.session_state.subject, st.session_state.chapter)
            st.session_state.selected_topics = st.multiselect("Select Topics to Cover", topics,
                                                              key="LECTUREPLANNERTopicsSelector",
                                                              default=st.session_state.selected_topics)

            if st.session_state.selected_topics:
                st.session_state.lecture_length = st.number_input(
                    "Lecture Length (minutes)",
                    min_value=15,
                    max_value=180,
                    value=st.session_state.lecture_length,
                    step=15,
                    key="LECTUREPLANNERLengthInput"
                )

                if st.button("Generate Presentation Structure", key="LECTUREPLANNERPresentationButtonCreator"):
                    with st.spinner("Generating presentation structure..."):
                        st.session_state.structure = generate_presentation_structure(
                            st.session_state.subject,
                            st.session_state.chapter,
                            st.session_state.selected_topics,
                            st.session_state.lecture_length
                        )
                        st.session_state.parsed_structure = parse_presentation_structure(st.session_state.structure, st.session_state.subject, st.session_state.chapter, st.session_state.selected_topics)
                        st.success("Presentation structure generated successfully!")
                        st.write(f"Generated {len(st.session_state.parsed_structure)} slides")

                if st.session_state.structure:
                    st.subheader("Generated Presentation Structure")
                    st.text_area("Presentation Outline", value=st.session_state.structure, height=400,
                                 key="generated_structure")

                    st.subheader("Review and Edit Presentation Structure")
                    edited_structure = review_and_edit_structure(st.session_state.parsed_structure)

                    if st.button("Update Structure"):
                        st.session_state.parsed_structure = edited_structure
                        st.success("Structure updated successfully!")
                        st.write(f"Updated structure now has {len(st.session_state.parsed_structure)} slides")

                    if st.button("Create PowerPoint Presentation"):
                        if st.session_state.parsed_structure:
                            with st.spinner("Creating PowerPoint presentation..."):
                                pptx_buffer = create_powerpoint(st.session_state.parsed_structure)

                            if pptx_buffer:
                                st.download_button(
                                    label="Download PowerPoint Presentation",
                                    data=pptx_buffer,
                                    file_name="lecture_presentation.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                )
                                st.success("PowerPoint presentation created successfully!")
                            else:
                                st.error("Failed to create PowerPoint presentation")
                        else:
                            st.error(
                                "No presentation structure available. Please generate or update the structure first.")